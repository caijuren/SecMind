#!/usr/bin/env python3
"""Build the Secmind Atomic Investigation Skills technical review deck.

Output: 2026-05-24-atomic-skills-demo.pptx (22 slides, 16:9 widescreen).

Run:  python3 build_pptx.py
"""
from __future__ import annotations

import datetime as _dt
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# -------------------------------------------------------------------------
# Palette / constants
# -------------------------------------------------------------------------
NAVY = RGBColor(0x1F, 0x2A, 0x44)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
GRAY_BG = RGBColor(0xF3, 0xF4, 0xF6)
GRAY_MID = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_TEXT = RGBColor(0x4B, 0x55, 0x63)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x11, 0x18, 0x27)

HEADER_TEXT = "Secmind · Atomic Investigation Skills"
DATE_STR = "2026-05-24"
OUTPUT_PATH = Path(
    "/Users/asoiso/workspace/Secmind/docs/slides/"
    "2026-05-24-atomic-skills-demo.pptx"
)

FONT_BODY = "Microsoft YaHei"
FONT_MONO = "Menlo"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# -------------------------------------------------------------------------
# Low-level helpers
# -------------------------------------------------------------------------
def _set_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _no_border(shape):
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, fill: RGBColor, line: bool = False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if not line:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = GRAY_MID
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = DARK_TEXT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font: str = FONT_BODY,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_header_footer(slide, page_no: int, total: int):
    # top header bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.45), NAVY)
    add_text(
        slide,
        Inches(0.4),
        Inches(0.08),
        Inches(10),
        Inches(0.35),
        HEADER_TEXT,
        size=12,
        color=WHITE,
        bold=True,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        Inches(10.5),
        Inches(0.08),
        Inches(2.6),
        Inches(0.35),
        DATE_STR,
        size=11,
        color=GRAY_MID,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # bottom footer
    add_rect(slide, Inches(0), Inches(7.25), SLIDE_W, Inches(0.25), GRAY_BG)
    add_text(
        slide,
        Inches(0.4),
        Inches(7.25),
        Inches(8),
        Inches(0.25),
        "secmind-investigator · P0 Technical Review",
        size=9,
        color=GRAY_TEXT,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        Inches(11.5),
        Inches(7.25),
        Inches(1.6),
        Inches(0.25),
        f"{page_no} / {total}",
        size=9,
        color=GRAY_TEXT,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_title(slide, title: str, subtitle: str | None = None):
    add_text(
        slide,
        Inches(0.5),
        Inches(0.65),
        Inches(12.3),
        Inches(0.7),
        title,
        size=30,
        bold=True,
        color=NAVY,
    )
    # accent underline bar
    add_rect(slide, Inches(0.5), Inches(1.35), Inches(0.9), Inches(0.06), BLUE)
    if subtitle:
        add_text(
            slide,
            Inches(0.5),
            Inches(1.45),
            Inches(12.3),
            Inches(0.45),
            subtitle,
            size=14,
            color=GRAY_TEXT,
        )


def add_bullets(slide, left, top, width, height, items, *, size=18, color=DARK_TEXT,
                line_spacing=1.25, bullet_color=BLUE):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet dot run
        bullet_run = p.add_run()
        bullet_run.text = "■ "
        bullet_run.font.name = FONT_BODY
        bullet_run.font.size = Pt(size)
        bullet_run.font.bold = True
        bullet_run.font.color.rgb = bullet_color
        # text run
        run = p.add_run()
        run.text = item
        run.font.name = FONT_BODY
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_code_block(slide, left, top, width, height, code: str, *, size: int = 13):
    add_rect(slide, left, top, width, height, RGBColor(0x0F, 0x17, 0x2A))
    tb = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), height - Inches(0.2),
    )
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = FONT_MONO
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)


def add_table(slide, left, top, width, height, rows, *,
              header_fill=NAVY, header_text=WHITE,
              body_text=DARK_TEXT, font_size=12, header_size=13,
              col_widths=None):
    """rows: list of list[str]; first row is header."""
    nrows = len(rows)
    ncols = len(rows[0])
    tbl_shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = tbl_shape.table
    # column widths
    if col_widths:
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = w
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                fc = header_text
                fs = header_size
                fb = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else GRAY_BG
                fc = body_text
                fs = font_size
                fb = False
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(0)
            tf.margin_right = Emu(0)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            # clear default run
            if p.runs:
                run = p.runs[0]
                run.text = str(val)
            else:
                run = p.add_run()
                run.text = str(val)
            run.font.name = FONT_BODY
            run.font.size = Pt(fs)
            run.font.bold = fb
            run.font.color.rgb = fc
    return tbl


def add_callout(slide, left, top, width, height, text, *, fill=BLUE, color=WHITE,
                size=14, bold=True):
    add_rect(slide, left, top, width, height, fill)
    add_text(
        slide, left, top, width, height, text,
        size=size, bold=bold, color=color,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )


# -------------------------------------------------------------------------
# Slide builders
# -------------------------------------------------------------------------
def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def slide_1_cover(prs, total):
    s = new_slide(prs)
    # full navy background
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    # accent bar
    add_rect(s, Inches(0.5), Inches(2.0), Inches(0.18), Inches(3.5), BLUE)
    add_text(
        s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2),
        "安全告警调查 · 原子能力实施计划",
        size=44, bold=True, color=WHITE,
    )
    add_text(
        s, Inches(0.9), Inches(3.2), Inches(11.5), Inches(0.7),
        "Atomic Investigation Capabilities — P0 技术评审",
        size=22, color=GRAY_MID,
    )
    add_text(
        s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.5),
        f"日期 · {DATE_STR}",
        size=18, color=BLUE, bold=True,
    )
    add_text(
        s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.4),
        "Secmind / secmind-investigator",
        size=14, color=GRAY_MID,
    )
    # no header/footer on cover; light page marker
    add_text(
        s, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3),
        f"1 / {total}",
        size=10, color=GRAY_MID, align=PP_ALIGN.RIGHT,
    )


def slide_2_agenda(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "议程", "本次评审的 8 个章节")
    items = [
        ("01", "背景与目标"),
        ("02", "三层架构"),
        ("03", "统一 Schema"),
        ("04", "数据源抽象"),
        ("05", "原子能力清单"),
        ("06", "错误与降级语义"),
        ("07", "端到端钓鱼场景"),
        ("08", "P0 验收 + 路线图"),
    ]
    # 4 columns x 2 rows
    box_w = Inches(2.95)
    box_h = Inches(2.1)
    gap_x = Inches(0.15)
    gap_y = Inches(0.3)
    start_x = Inches(0.5)
    start_y = Inches(2.3)
    for i, (idx, title) in enumerate(items):
        row = i // 4
        col = i % 4
        x = start_x + (box_w + gap_x) * col
        y = start_y + (box_h + gap_y) * row
        add_rect(s, x, y, box_w, box_h, GRAY_BG)
        add_rect(s, x, y, Inches(0.1), box_h, BLUE)
        add_text(s, x + Inches(0.3), y + Inches(0.25), box_w - Inches(0.4),
                 Inches(0.6), idx, size=28, bold=True, color=BLUE)
        add_text(s, x + Inches(0.3), y + Inches(1.0), box_w - Inches(0.4),
                 Inches(1.0), title, size=20, bold=True, color=NAVY)


def slide_3_background(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "背景与目标", "为什么需要把调查能力做成可被 LLM Agent 调用的积木")

    # left: pain points
    add_rect(s, Inches(0.5), Inches(2.1), Inches(6.0), Inches(2.0), GRAY_BG)
    add_rect(s, Inches(0.5), Inches(2.1), Inches(0.1), Inches(2.0), RED)
    add_text(s, Inches(0.75), Inches(2.2), Inches(5.7), Inches(0.5),
             "痛点", size=18, bold=True, color=RED)
    add_bullets(
        s, Inches(0.75), Inches(2.7), Inches(5.6), Inches(1.4),
        [
            "告警调查严重依赖人工经验,新人上手周期长",
            "跨 IDP / EDR / SIEM / 邮件 等系统反复拼接",
            "调查流程难以复用、难以审计、难以度量",
        ],
        size=14, line_spacing=1.3,
    )

    # right: goal
    add_rect(s, Inches(6.85), Inches(2.1), Inches(6.0), Inches(2.0), GRAY_BG)
    add_rect(s, Inches(6.85), Inches(2.1), Inches(0.1), Inches(2.0), GREEN)
    add_text(s, Inches(7.1), Inches(2.2), Inches(5.7), Inches(0.5),
             "目标", size=18, bold=True, color=GREEN)
    add_bullets(
        s, Inches(7.1), Inches(2.7), Inches(5.6), Inches(1.4),
        [
            "把调查动作拆解成可独立调用的原子能力",
            "LLM Agent / 上层编排器自由组合这些积木",
            "Schema 统一、证据可回溯、降级可观测",
        ],
        size=14, line_spacing=1.3,
    )

    # bottom: deliverable + phases
    add_rect(s, Inches(0.5), Inches(4.4), Inches(12.35), Inches(2.4), GRAY_BG)
    add_text(s, Inches(0.75), Inches(4.55), Inches(11.8), Inches(0.5),
             "产物与节奏", size=18, bold=True, color=NAVY)
    add_text(s, Inches(0.75), Inches(5.05), Inches(11.8), Inches(0.45),
             "Python 库:secmind-investigator(Python 3.11+ / Pydantic v2 / asyncio)",
             size=14, color=DARK_TEXT)
    # phases
    px = [Inches(0.75), Inches(4.95), Inches(9.15)]
    labels = [
        ("P0 · 基础底座", "Schema + 9 类数据源 + 25 原子能力", BLUE),
        ("P1 · 场景扩展", "+6 数据源 + 云/横移/外泄/特权能力", GREEN),
        ("P2 · 深度专题", "沙箱回放 / 漏洞链 / 域名全景", NAVY),
    ]
    for x, (h, sub, col) in zip(px, labels):
        add_rect(s, x, Inches(5.7), Inches(3.9), Inches(1.0), WHITE, line=True)
        add_rect(s, x, Inches(5.7), Inches(0.08), Inches(1.0), col)
        add_text(s, x + Inches(0.18), Inches(5.78), Inches(3.7), Inches(0.4),
                 h, size=14, bold=True, color=col)
        add_text(s, x + Inches(0.18), Inches(6.18), Inches(3.7), Inches(0.5),
                 sub, size=12, color=GRAY_TEXT)


def slide_4_three_layers(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "三层架构总览", "Core / Adapter / Capability — 单向依赖")

    layers = [
        ("Capability 层", "每个能力一个独立类,单一职责,不互相调用",
         "investigate.entity.user.profile / investigate.verdict.tp_fp ...", BLUE),
        ("DataSource Adapter 层", "每类数据源(DS-XXX)一份接口,可挂多家厂商实现",
         "DS-IDP / DS-EDR / DS-SIEM / DS-EMAIL / DS-TI ...", GREEN),
        ("Core 层", "Schema / AtomicCapability ABC / Registry / 错误与降级 / Observability",
         "CapabilityResult · EvidenceRef · PIIMarker", NAVY),
    ]
    top = Inches(2.0)
    h = Inches(1.4)
    gap = Inches(0.15)
    for i, (name, desc, ex, col) in enumerate(layers):
        y = top + (h + gap) * i
        add_rect(s, Inches(0.5), y, Inches(12.35), h, GRAY_BG)
        add_rect(s, Inches(0.5), y, Inches(0.18), h, col)
        add_text(s, Inches(0.85), y + Inches(0.15), Inches(4.0), Inches(0.5),
                 name, size=20, bold=True, color=col)
        add_text(s, Inches(0.85), y + Inches(0.7), Inches(11.5), Inches(0.4),
                 desc, size=14, color=DARK_TEXT)
        add_text(s, Inches(0.85), y + Inches(1.05), Inches(11.5), Inches(0.35),
                 ex, size=12, color=GRAY_TEXT, font=FONT_MONO)

    # bottom note
    add_callout(
        s, Inches(0.5), Inches(6.5), Inches(12.35), Inches(0.55),
        "本仓库提供积木,编排由 LLM Agent / 上层规则引擎负责",
        fill=NAVY, color=WHITE, size=14, bold=True,
    )


def slide_5_stack(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "技术栈与工程基线", "可重复、可观测、可门禁")

    # left: tech list
    add_text(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(0.4),
             "运行时与依赖", size=16, bold=True, color=NAVY)
    add_bullets(
        s, Inches(0.5), Inches(2.45), Inches(6.2), Inches(3.5),
        [
            "Python 3.11+ / Pydantic v2 / asyncio",
            "pytest + pytest-asyncio + pytest-cov + hypothesis",
            "structlog + OpenTelemetry(traces / metrics)",
            "ruff(lint+format) + mypy --strict",
            "uv:依赖解析 / venv / 锁文件",
        ],
        size=14, line_spacing=1.5,
    )

    # right: skeleton table
    add_text(s, Inches(6.9), Inches(2.0), Inches(6.0), Inches(0.4),
             "项目骨架与门禁", size=16, bold=True, color=NAVY)
    rows = [
        ["文件 / 配置", "作用"],
        ["pyproject.toml", "依赖、构建、ruff / mypy / pytest 配置"],
        ["Makefile", "install / lint / test / typecheck / coverage"],
        ["conftest.py", "公共 fixtures(InMemoryDataSource 等)"],
        [".github/workflows/", "CI:lint + typecheck + test"],
        ["coverage", "pytest --cov-fail-under=85(硬门禁)"],
    ]
    add_table(
        s, Inches(6.9), Inches(2.5), Inches(6.0), Inches(3.5), rows,
        font_size=12, header_size=13,
        col_widths=[Inches(2.0), Inches(4.0)],
    )

    add_callout(
        s, Inches(0.5), Inches(6.5), Inches(12.35), Inches(0.55),
        "make lint && make typecheck && make test —— PR 合入硬门禁",
        fill=BLUE, color=WHITE, size=13, bold=True,
    )


def slide_6_capability_result(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "统一 Schema:CapabilityResult", "所有原子能力共享同一信封")

    rows = [
        ["字段", "类型", "说明"],
        ["payload", "BaseModel", "能力强类型输出(每个能力自定义)"],
        ["confidence", "float [0,1]", "置信度,LLM 据此判断是否补查"],
        ["partial", "bool", "是否为降级 / 部分结果"],
        ["partial_reasons", "list[str]", "降级原因(可枚举),便于 LLM 决策"],
        ["evidence_refs", "list[EvidenceRef]", "证据回溯锚点(source/query/record_id)"],
        ["pii_fields", "list[PIIMarker]", "标注敏感字段,下游 Redactor 按角色裁剪"],
        ["duration_ms", "int", "执行耗时,接入 SLO"],
        ["cost_units", "float | None", "调用代价(API 计费 / 配额)"],
    ]
    add_table(
        s, Inches(0.5), Inches(2.0), Inches(12.35), Inches(4.4), rows,
        font_size=12, header_size=13,
        col_widths=[Inches(2.4), Inches(2.4), Inches(7.55)],
    )

    add_callout(
        s, Inches(0.5), Inches(6.55), Inches(12.35), Inches(0.5),
        "统一信封 → LLM 无需为每个能力学新结构,只学一次",
        fill=GREEN, color=WHITE, size=13,
    )


def slide_7_schema_types(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "Schema 配套类型", "EvidenceRef / PIIMarker — 可审计 + 可裁剪")

    add_text(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(0.4),
             "EvidenceRef · 可审计回溯", size=16, bold=True, color=BLUE)
    add_bullets(
        s, Inches(0.5), Inches(2.45), Inches(6.0), Inches(1.6),
        [
            "source:数据源 id(如 ds.idp.entra)",
            "query:执行的查询语句 / API 名",
            "record_id:返回记录的主键",
            "captured_at:抓取时间戳(UTC)",
        ],
        size=13, line_spacing=1.4,
    )

    add_text(s, Inches(0.5), Inches(4.2), Inches(6.0), Inches(0.4),
             "PIIMarker · 敏感字段裁剪", size=16, bold=True, color=RED)
    add_bullets(
        s, Inches(0.5), Inches(4.65), Inches(6.0), Inches(1.6),
        [
            "json_path:目标字段的 JSON Path",
            "classification:pii / secret / internal / regulated",
            "下游 Redactor 按访问角色统一脱敏",
        ],
        size=13, line_spacing=1.4,
    )

    code = (
        "class EvidenceRef(BaseModel):\n"
        "    source: str\n"
        "    query: str\n"
        "    record_id: str | None = None\n"
        "    captured_at: datetime\n"
        "\n"
        "class PIIMarker(BaseModel):\n"
        "    json_path: str\n"
        "    classification: Literal[\n"
        "        'pii', 'secret', 'internal', 'regulated'\n"
        "    ]"
    )
    add_code_block(
        s, Inches(6.9), Inches(2.0), Inches(6.0), Inches(4.6), code, size=13,
    )


def slide_8_atomic_capability(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "核心抽象:AtomicCapability", "单一职责 · 不互调 · 强 Schema")

    add_bullets(
        s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(3.0),
        [
            "每个能力声明 namespace + input_model + output_payload_model",
            "execute() 由子类实现,只做一件事",
            "run() 统一封装:输入校验 / 计时 / 异常映射 / 日志 trace",
            "禁止能力间互相调用,组合由 LLM Agent 完成",
        ],
        size=14, line_spacing=1.5,
    )

    code = (
        "class AtomicCapability(ABC, Generic[I, P]):\n"
        "    namespace: ClassVar[str]\n"
        "    input_model: ClassVar[type[BaseModel]]\n"
        "    output_payload_model: ClassVar[type[BaseModel]]\n"
        "\n"
        "    async def run(self, raw: dict) -> CapabilityResult[P]:\n"
        "        inp = self.input_model.model_validate(raw)\n"
        "        t0 = time.monotonic()\n"
        "        try:\n"
        "            payload = await self.execute(inp)\n"
        "            return CapabilityResult(payload=payload, ...)\n"
        "        except DataSourceUnavailable as e:\n"
        "            return self._partial(reasons=[str(e)])\n"
        "\n"
        "    @abstractmethod\n"
        "    async def execute(self, inp: I) -> P: ..."
    )
    add_code_block(
        s, Inches(6.9), Inches(2.0), Inches(6.0), Inches(4.8), code, size=12,
    )


def slide_9_registry(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "命名空间路由:CapabilityRegistry", "稳定 namespace + 装饰器注册")

    add_bullets(
        s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(3.0),
        [
            "namespace 形如 investigate.<category>.<name>",
            "@register 装饰器:导入即注册",
            "Registry.resolve(ns) 返回单例能力实例",
            "前缀枚举(Categories)避免拼写漂移",
        ],
        size=14, line_spacing=1.5,
    )

    rows = [
        ["类别", "示例 namespace"],
        ["entity", "investigate.entity.user.profile"],
        ["entity", "investigate.entity.host.profile"],
        ["activity", "investigate.activity.login"],
        ["intel", "investigate.intel.ip.reputation"],
        ["relation", "investigate.relation.entity_graph"],
        ["verdict", "investigate.verdict.tp_fp"],
        ["meta", "investigate.meta.next_step"],
    ]
    add_table(
        s, Inches(6.9), Inches(2.0), Inches(6.0), Inches(4.6), rows,
        font_size=12, header_size=13,
        col_widths=[Inches(1.5), Inches(4.5)],
    )


def slide_10_errors(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "错误与降级语义", "三种异常分级 · LLM 可读")

    rows = [
        ["异常", "行为", "触发场景"],
        [
            "DataSourceUnavailable",
            "run() 转 partial=True · confidence=0.0 · 不抛",
            "API 429 / 5xx / 网络抖动 / 健康检查失败",
        ],
        [
            "PartialResult",
            "携带部分 payload + reasons + 较低 confidence",
            "部分字段缺失 / 只命中部分数据源 / 超时降级",
        ],
        [
            "CapabilityFatal",
            "抛出 · 由上层编排器决定中断或重试",
            "输入校验失败 / 配置缺失 / 编程错误",
        ],
    ]
    add_table(
        s, Inches(0.5), Inches(2.0), Inches(12.35), Inches(3.2), rows,
        font_size=12, header_size=13,
        col_widths=[Inches(2.6), Inches(4.7), Inches(5.05)],
    )

    # bottom illustration
    add_rect(s, Inches(0.5), Inches(5.5), Inches(12.35), Inches(1.4), GRAY_BG)
    add_text(s, Inches(0.75), Inches(5.6), Inches(11.85), Inches(0.4),
             "降级路径示意", size=14, bold=True, color=NAVY)
    add_text(s, Inches(0.75), Inches(6.0), Inches(11.85), Inches(0.4),
             "DataSourceUnavailable → CapabilityResult(partial=True, partial_reasons=['ds.idp.entra timeout'])",
             size=12, color=GRAY_TEXT, font=FONT_MONO)
    add_text(s, Inches(0.75), Inches(6.4), Inches(11.85), Inches(0.4),
             "LLM 读 partial_reasons → 决策:跳过 / 切换数据源 / 补查另一个能力",
             size=12, color=NAVY, bold=True)


def slide_11_observability(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "可观测性接入", "每次 run() 自动出日志 + Metric + Trace")

    # logs
    add_text(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(0.4),
             "Structured Log(structlog)", size=16, bold=True, color=BLUE)
    code_log = (
        "{\n"
        '  "event": "capability.run",\n'
        '  "capability_ns": "investigate.entity.user.profile",\n'
        '  "trace_id": "9c3f...",\n'
        '  "duration_ms": 142,\n'
        '  "partial": false,\n'
        '  "confidence": 0.92\n'
        "}"
    )
    add_code_block(s, Inches(0.5), Inches(2.45), Inches(6.0), Inches(2.6),
                   code_log, size=12)

    # metrics
    add_text(s, Inches(6.9), Inches(2.0), Inches(6.0), Inches(0.4),
             "OpenTelemetry Metric", size=16, bold=True, color=GREEN)
    rows = [
        ["指标名", "类型", "标签"],
        ["capability.run.total", "Counter", "ns / partial"],
        ["capability.partial.total", "Counter", "ns / reason"],
        ["capability.run.duration_ms", "Histogram", "ns"],
        ["datasource.call.total", "Counter", "ds_id / status"],
    ]
    add_table(s, Inches(6.9), Inches(2.45), Inches(6.0), Inches(2.6), rows,
              font_size=11, header_size=12,
              col_widths=[Inches(2.5), Inches(1.5), Inches(2.0)])

    add_callout(
        s, Inches(0.5), Inches(5.4), Inches(12.35), Inches(1.4),
        "为 SLO 监控、故障定位、成本审计打基础\n"
        "(P95 延迟、partial 率、按能力 / 数据源拆分)",
        fill=NAVY, color=WHITE, size=14, bold=True,
    )


def slide_12_datasources(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "数据源抽象:DS-XXX 一览", "P0 必备 9 类 + P1 预留 6 类")

    rows_p0 = [
        ["缩写", "含义", "典型厂商"],
        ["DS-IDP", "身份与认证", "Entra ID / Okta"],
        ["DS-HR", "人事系统", "Workday / 飞书"],
        ["DS-CMDB", "资产配置", "ServiceNow / Lansweeper"],
        ["DS-EDR", "终端检测响应", "CrowdStrike / SentinelOne"],
        ["DS-SIEM", "日志与告警", "Splunk / Sentinel / Chronicle"],
        ["DS-EMAIL", "邮件安全", "M365 / Proofpoint"],
        ["DS-FW", "网络防火墙", "Palo Alto / Fortinet"],
        ["DS-TI", "威胁情报", "VT / OTX / 商业 TI"],
        ["DS-KB", "知识库", "Confluence / Runbook"],
    ]
    add_text(s, Inches(0.5), Inches(1.95), Inches(6.0), Inches(0.4),
             "P0 必需(9 类)", size=14, bold=True, color=BLUE)
    add_table(s, Inches(0.5), Inches(2.4), Inches(6.0), Inches(4.4), rows_p0,
              font_size=11, header_size=12,
              col_widths=[Inches(1.2), Inches(2.0), Inches(2.8)])

    rows_p1 = [
        ["缩写", "含义"],
        ["DS-NDR", "网络检测响应"],
        ["DS-CLOUD", "云审计日志"],
        ["DS-SANDBOX", "沙箱动态分析"],
        ["DS-PAM", "特权访问管理"],
        ["DS-VULN", "漏洞扫描"],
        ["DS-WHOIS", "域名注册信息"],
    ]
    add_text(s, Inches(6.9), Inches(1.95), Inches(6.0), Inches(0.4),
             "P1 预留(6 类)", size=14, bold=True, color=GREEN)
    add_table(s, Inches(6.9), Inches(2.4), Inches(6.0), Inches(3.5), rows_p1,
              font_size=11, header_size=12,
              col_widths=[Inches(2.0), Inches(4.0)])

    add_callout(
        s, Inches(6.9), Inches(6.05), Inches(6.0), Inches(0.7),
        "圆点约定:● 必需  ○ 可选",
        fill=NAVY, color=WHITE, size=13,
    )


def slide_13_datasource_abc(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "DataSource ABC · 限速 · Mock", "可替换的厂商接入点")

    add_bullets(
        s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(3.0),
        [
            "DataSource ABC 强制实现 id / healthcheck() / close()",
            "@rate_limited(rps=...) 装饰器:简单令牌桶",
            "InMemoryDataSource 供单测 + Demo 复用",
            "厂商实现:子类化 ABC + 注册到 DataSourceRegistry",
        ],
        size=14, line_spacing=1.5,
    )

    code = (
        "class DataSource(ABC):\n"
        "    id: ClassVar[str]\n"
        "\n"
        "    @abstractmethod\n"
        "    async def healthcheck(self) -> bool: ...\n"
        "\n"
        "    @abstractmethod\n"
        "    async def close(self) -> None: ...\n"
        "\n"
        "@rate_limited(rps=20)\n"
        "class EntraIDPSource(DataSource):\n"
        "    id = 'ds.idp.entra'\n"
        "    async def lookup_user(self, upn): ..."
    )
    add_code_block(s, Inches(6.9), Inches(2.0), Inches(6.0), Inches(4.6),
                   code, size=12)


def slide_14_capability_inventory(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "P0 原子能力清单(类别视图)", "6 大类别 · 共 25 个能力")

    rows = [
        ["类别", "能力数", "示例"],
        ["Entity 实体", "5", "A1 user / A2 host / A4 ip / A5 domain / A7 file"],
        ["Activity 活动", "4", "B1 login / B3 process / B4 network / B8 email"],
        ["Relation 关联", "3", "C1 cooccurrence / C2 entity_graph / C3 process_lineage"],
        ["Context 上下文", "3", "D1 asset_business / D2 user_identity / D5 baseline_self"],
        ["Intel 情报", "4", "E1 ip_rep / E2 domain_url_rep / E3 hash_rep / E7 ttp_mitre"],
        ["Verdict + Meta", "6", "F1 TP/FP · F2 impact · F4 ioc · F8 response · G3 next · G4 kb"],
    ]
    add_table(
        s, Inches(0.5), Inches(2.0), Inches(12.35), Inches(4.4), rows,
        font_size=12, header_size=13,
        col_widths=[Inches(2.4), Inches(1.2), Inches(8.75)],
    )

    add_callout(
        s, Inches(0.5), Inches(6.55), Inches(12.35), Inches(0.5),
        "25 个能力全部进 Registry,LLM Agent 通过 namespace 直接调用",
        fill=BLUE, color=WHITE, size=13,
    )


def slide_15_a1_user_profile(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "重点能力示例:A1 用户/账户画像",
              "investigate.entity.user.profile")

    # left meta
    rows = [
        ["项", "值"],
        ["namespace", "investigate.entity.user.profile"],
        ["输入", "principal:str(邮箱 / UPN)"],
        ["数据源", "DS-IDP●  DS-HR●  DS-CMDB○  DS-EMAIL○"],
        ["典型耗时", "< 800 ms (P95,Mock 数据源)"],
    ]
    add_table(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(2.6), rows,
              font_size=12, header_size=13,
              col_widths=[Inches(1.4), Inches(4.6)])

    # right payload fields
    add_text(s, Inches(6.9), Inches(2.0), Inches(6.0), Inches(0.4),
             "Payload 关键字段", size=14, bold=True, color=NAVY)
    add_bullets(
        s, Inches(6.9), Inches(2.45), Inches(6.0), Inches(3.0),
        [
            "identity:UPN / 状态 / MFA / 风险等级",
            "hr:部门 / 岗位 / 直属上级 / 在/离职态",
            "device_links:绑定的设备 id 列表",
            "mailbox:邮箱归属 / 转发规则",
            "risk_signals:近 30 天高危登录 / 异地告警",
        ],
        size=13, line_spacing=1.4,
    )

    add_callout(
        s, Inches(0.5), Inches(6.05), Inches(12.35), Inches(0.8),
        "LLM 据此判断:是否高敏感用户 / 是否启用 MFA / 是否离职态",
        fill=NAVY, color=WHITE, size=14, bold=True,
    )


def slide_16_f1_verdict(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "重点能力示例:F1 TP/FP 判定",
              "investigate.verdict.tp_fp")

    rows = [
        ["项", "值"],
        ["namespace", "investigate.verdict.tp_fp"],
        ["输入", "evidence_set: list[CapabilityResult]"],
        ["输出 verdict", "TP / FP / BP(Benign Positive) / INC(Inconclusive)"],
        ["输出附加", "rationale · missing_evidence · confidence"],
    ]
    add_table(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(2.6), rows,
              font_size=12, header_size=13,
              col_widths=[Inches(1.6), Inches(4.4)])

    # right: loop diagram
    add_text(s, Inches(6.9), Inches(2.0), Inches(6.0), Inches(0.4),
             "判定 → 不确定性 → 补查 闭环", size=14, bold=True, color=NAVY)
    # three boxes
    box_w = Inches(1.85)
    box_h = Inches(0.9)
    yb = Inches(2.6)
    xs = [Inches(6.9), Inches(8.95), Inches(11.0)]
    titles = [("F1\nTP/FP", BLUE), ("missing\n_evidence", RED), ("G3\nnext_step", GREEN)]
    for x, (t, c) in zip(xs, titles):
        add_rect(s, x, yb, box_w, box_h, c)
        add_text(s, x, yb, box_w, box_h, t, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # arrows (simple textual arrows)
    add_text(s, Inches(8.78), yb, Inches(0.25), box_h, "→",
             size=22, bold=True, color=GRAY_TEXT, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(10.83), yb, Inches(0.25), box_h, "→",
             size=22, bold=True, color=GRAY_TEXT, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    add_bullets(
        s, Inches(6.9), Inches(3.8), Inches(6.0), Inches(2.5),
        [
            "F1 给出判定 + 缺失证据列表",
            "G3 next_step 把缺失证据映射成下一步能力调用",
            "形成可解释、可追问的调查闭环",
        ],
        size=12, line_spacing=1.4,
    )

    add_callout(
        s, Inches(0.5), Inches(6.05), Inches(12.35), Inches(0.8),
        "Verdict ≠ 终结,Verdict + missing_evidence + next_step = 自驱调查",
        fill=NAVY, color=WHITE, size=13, bold=True,
    )


def slide_17_phishing_flow(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "端到端钓鱼场景演示", "Alert → 6 个原子能力 → containment 建议")

    steps = [
        ("01", "A1 UserProfile", "查收件人画像 / MFA / 离职态", BLUE),
        ("02", "B8 EmailActivity", "调取告警邮件元数据 / 链接", BLUE),
        ("03", "E2 URL Reputation", "查链接信誉 / 域名年龄 / 黑名单", GREEN),
        ("04", "F1 TP/FP Verdict", "verdict = TP · confidence 0.91", RED),
        ("05", "F4 IOC Extract", "≥3 IOCs(URL / 域名 / 哈希)", NAVY),
        ("06", "F8 Response", "containment 建议 + Runbook 引用", NAVY),
    ]
    box_w = Inches(2.0)
    box_h = Inches(2.3)
    gap = Inches(0.07)
    y = Inches(2.2)
    x0 = Inches(0.5)
    for i, (idx, name, desc, col) in enumerate(steps):
        x = x0 + (box_w + gap) * i
        add_rect(s, x, y, box_w, box_h, GRAY_BG)
        add_rect(s, x, y, box_w, Inches(0.4), col)
        add_text(s, x, y, box_w, Inches(0.4), idx,
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.1), y + Inches(0.5),
                 box_w - Inches(0.2), Inches(0.6),
                 name, size=13, bold=True, color=NAVY)
        add_text(s, x + Inches(0.1), y + Inches(1.05),
                 box_w - Inches(0.2), Inches(1.2),
                 desc, size=11, color=GRAY_TEXT)

    # alert entry
    add_text(s, Inches(0.5), Inches(1.85), Inches(12.35), Inches(0.3),
             "Alert(钓鱼链接)→",
             size=14, bold=True, color=RED)

    add_callout(
        s, Inches(0.5), Inches(5.0), Inches(12.35), Inches(1.0),
        "6 个原子能力,全部 Mock 数据源,端到端 < 2 秒完成",
        fill=GREEN, color=WHITE, size=18, bold=True,
    )
    # closing line
    add_text(
        s, Inches(0.5), Inches(6.2), Inches(12.35), Inches(0.6),
        "演示意义:验证 Schema / Registry / 错误降级 / Observability 链路通畅",
        size=13, color=GRAY_TEXT, align=PP_ALIGN.CENTER,
    )


def slide_18_phishing_code(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "钓鱼场景代码片段",
              "test_phishing_orchestration.py(精简)")

    code = (
        "async def test_phishing_e2e(registry, mock_sources):\n"
        "    alert = {'recipient': 'alice@corp', 'url': 'http://evil.tld'}\n"
        "\n"
        "    user   = await registry.resolve(\n"
        "        'investigate.entity.user.profile'\n"
        "    ).run({'principal': alert['recipient']})\n"
        "\n"
        "    email  = await registry.resolve(\n"
        "        'investigate.activity.email'\n"
        "    ).run({'alert_id': 'A-001'})\n"
        "\n"
        "    rep    = await registry.resolve(\n"
        "        'investigate.intel.domain_url.reputation'\n"
        "    ).run({'url': alert['url']})\n"
        "\n"
        "    verdict = await registry.resolve(\n"
        "        'investigate.verdict.tp_fp'\n"
        "    ).run({'evidence_set': [user, email, rep]})\n"
        "\n"
        "    assert verdict.payload.verdict == 'TP'\n"
        "    iocs    = await registry.resolve('investigate.verdict.ioc').run(...)\n"
        "    plan    = await registry.resolve('investigate.verdict.response').run(...)\n"
        "    assert len(iocs.payload.indicators) >= 3"
    )
    add_code_block(s, Inches(0.5), Inches(2.0), Inches(12.35), Inches(4.6),
                   code, size=12)

    add_callout(
        s, Inches(0.5), Inches(6.7), Inches(12.35), Inches(0.45),
        "重点:每一步都是 await cap.run(payload) 的统一调用",
        fill=BLUE, color=WHITE, size=13,
    )


def slide_19_p0_tasks(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "P0 任务总览(43 个 Task)", "Foundation → Adapter → Capability → 验收")

    rows = [
        ["阶段", "Task 数", "内容要点"],
        ["Foundation 基础底座", "6", "Schema · ABC · Errors · Registry · Observability · Skeleton"],
        ["DataSource Adapter", "9", "9 类 P0 数据源 + Mock 实现 + Healthcheck"],
        ["Capability 能力实现", "25", "25 个原子能力,逐个 PR,带单测"],
        ["Integration 集成测试", "1", "钓鱼场景 E2E + 编排示例"],
        ["Docs 文档", "1", "README · 能力索引 · 使用 Cookbook"],
        ["验收门禁", "1", "lint + typecheck + coverage ≥ 85% + registry ≥ 25"],
    ]
    add_table(s, Inches(0.5), Inches(2.0), Inches(12.35), Inches(3.6), rows,
              font_size=12, header_size=13,
              col_widths=[Inches(3.0), Inches(1.2), Inches(8.15)])

    # progress bar
    add_text(s, Inches(0.5), Inches(5.85), Inches(12.35), Inches(0.4),
             "当前进度", size=14, bold=True, color=NAVY)
    # background bar
    bar_x = Inches(0.5)
    bar_y = Inches(6.25)
    bar_w = Inches(12.35)
    bar_h = Inches(0.35)
    add_rect(s, bar_x, bar_y, bar_w, bar_h, GRAY_MID)
    # done: 3 of 43 ≈ 7%
    done_w = Inches(12.35 * 3 / 43)
    add_rect(s, bar_x, bar_y, done_w, bar_h, GREEN)
    add_text(s, bar_x, bar_y, bar_w, bar_h,
             "3 / 43 完成(Schema + EvidenceRef + PIIMarker 已提交)",
             size=12, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_callout(
        s, Inches(0.5), Inches(6.75), Inches(12.35), Inches(0.4),
        "门禁命令:make lint && make test && --cov-fail-under=85 && registry ≥ 25",
        fill=NAVY, color=WHITE, size=11,
    )


def slide_20_p1_roadmap(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "P1 路线图(Milestone-level)", "覆盖云 / 横移 / 外泄 / 特权")

    # new datasources
    rows_ds = [
        ["新增数据源", "估时"],
        ["DS-NDR(网络检测响应)", "1-2 d"],
        ["DS-CLOUD(云审计:AWS / Azure / GCP)", "2 d"],
        ["DS-DLP(数据防泄漏)", "1-2 d"],
        ["DS-CASB(SaaS 治理)", "1-2 d"],
        ["DS-DAM(数据库访问监控)", "1-2 d"],
        ["DS-PAM(特权访问管理)", "1-2 d"],
    ]
    add_text(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(0.4),
             "6 个新数据源", size=14, bold=True, color=BLUE)
    add_table(s, Inches(0.5), Inches(2.45), Inches(6.0), Inches(3.5), rows_ds,
              font_size=11, header_size=12,
              col_widths=[Inches(4.4), Inches(1.6)])

    # new capabilities
    rows_cap = [
        ["新增能力主题"],
        ["云审计调查(CloudTrail / Activity Log)"],
        ["横向移动检测(Lateral Movement)"],
        ["数据外泄关联(DLP + CASB + EDR)"],
        ["特权滥用分析(PAM + IDP 高权操作)"],
    ]
    add_text(s, Inches(6.9), Inches(2.0), Inches(6.0), Inches(0.4),
             "4 个能力主题", size=14, bold=True, color=GREEN)
    add_table(s, Inches(6.9), Inches(2.45), Inches(6.0), Inches(2.6), rows_cap,
              font_size=12, header_size=13,
              col_widths=[Inches(6.0)])

    add_callout(
        s, Inches(6.9), Inches(5.2), Inches(6.0), Inches(0.75),
        "验收:覆盖 4 类新场景的集成测试",
        fill=NAVY, color=WHITE, size=13, bold=True,
    )


def slide_21_p2_roadmap(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "P2 路线图", "沙箱回放 · 漏洞利用链 · 域名注册全景")

    items = [
        ("DS-SANDBOX", "沙箱动态分析 · 文件 / URL 引爆", BLUE,
         "回放可疑样本行为 · 提取动态 IOC"),
        ("DS-VULN", "漏洞扫描 · CVE / 资产暴露面", GREEN,
         "把告警资产与 CVE 关联 · 推断利用链"),
        ("DS-WHOIS", "域名注册 · 历史解析 · 相关域", NAVY,
         "钓鱼 / 仿冒域全景画像 · 注册人聚类"),
    ]
    box_h = Inches(1.4)
    gap = Inches(0.2)
    top = Inches(2.0)
    for i, (name, sub, col, detail) in enumerate(items):
        y = top + (box_h + gap) * i
        add_rect(s, Inches(0.5), y, Inches(12.35), box_h, GRAY_BG)
        add_rect(s, Inches(0.5), y, Inches(0.18), box_h, col)
        add_text(s, Inches(0.85), y + Inches(0.15), Inches(3.5), Inches(0.5),
                 name, size=20, bold=True, color=col)
        add_text(s, Inches(0.85), y + Inches(0.7), Inches(12), Inches(0.35),
                 sub, size=13, color=DARK_TEXT)
        add_text(s, Inches(0.85), y + Inches(1.0), Inches(12), Inches(0.35),
                 detail, size=12, color=GRAY_TEXT)

    add_callout(
        s, Inches(0.5), Inches(6.7), Inches(12.35), Inches(0.45),
        "P2 聚焦深度专题,而非铺面积;每个专题独立验收",
        fill=NAVY, color=WHITE, size=13,
    )


def slide_22_qa(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "Q&A · 下一步", "当前进展与即将启动的工作")

    # done column
    add_rect(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(3.8), GRAY_BG)
    add_rect(s, Inches(0.5), Inches(2.0), Inches(0.1), Inches(3.8), GREEN)
    add_text(s, Inches(0.75), Inches(2.1), Inches(5.6), Inches(0.5),
             "已完成", size=18, bold=True, color=GREEN)
    add_bullets(
        s, Inches(0.75), Inches(2.6), Inches(5.6), Inches(3.0),
        [
            "Task 1:项目骨架(pyproject / Makefile / conftest)",
            "Task 2:make install / test / lint 全通",
            "Task 3:CapabilityResult / EvidenceRef / PIIMarker Schema",
            "覆盖率 ≥ 85%,ruff + mypy strict 通过",
        ],
        size=13, line_spacing=1.4,
    )

    # next column
    add_rect(s, Inches(6.9), Inches(2.0), Inches(6.0), Inches(3.8), GRAY_BG)
    add_rect(s, Inches(6.9), Inches(2.0), Inches(0.1), Inches(3.8), BLUE)
    add_text(s, Inches(7.15), Inches(2.1), Inches(5.6), Inches(0.5),
             "下一步", size=18, bold=True, color=BLUE)
    add_bullets(
        s, Inches(7.15), Inches(2.6), Inches(5.6), Inches(3.0),
        [
            "Task 4:Errors —— 三种异常 + run() 降级映射",
            "Task 5:CapabilityRegistry + 装饰器注册",
            "Task 6:Observability —— structlog + OTel",
            "之后并行推 9 个 DS Adapter + 25 个能力",
        ],
        size=13, line_spacing=1.4,
    )

    add_callout(
        s, Inches(0.5), Inches(6.1), Inches(12.35), Inches(0.8),
        "讨论入口:tech@ouraca.ai · Secmind / secmind-investigator",
        fill=NAVY, color=WHITE, size=16, bold=True,
    )


# -------------------------------------------------------------------------
# Driver
# -------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 22
    slide_1_cover(prs, total)
    slide_2_agenda(prs, 2, total)
    slide_3_background(prs, 3, total)
    slide_4_three_layers(prs, 4, total)
    slide_5_stack(prs, 5, total)
    slide_6_capability_result(prs, 6, total)
    slide_7_schema_types(prs, 7, total)
    slide_8_atomic_capability(prs, 8, total)
    slide_9_registry(prs, 9, total)
    slide_10_errors(prs, 10, total)
    slide_11_observability(prs, 11, total)
    slide_12_datasources(prs, 12, total)
    slide_13_datasource_abc(prs, 13, total)
    slide_14_capability_inventory(prs, 14, total)
    slide_15_a1_user_profile(prs, 15, total)
    slide_16_f1_verdict(prs, 16, total)
    slide_17_phishing_flow(prs, 17, total)
    slide_18_phishing_code(prs, 18, total)
    slide_19_p0_tasks(prs, 19, total)
    slide_20_p1_roadmap(prs, 20, total)
    slide_21_p2_roadmap(prs, 21, total)
    slide_22_qa(prs, 22, total)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH, len(prs.slides)


if __name__ == "__main__":
    path, n = build()
    size_kb = path.stat().st_size / 1024
    print(f"OK: wrote {path} · slides={n} · size={size_kb:.1f} KB")
