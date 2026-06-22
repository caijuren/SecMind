#!/usr/bin/env python3
"""Build the Secmind Atomic Investigation Skills SPEC design-review deck.

Output: 2026-05-24-atomic-skills-spec.pptx (24 slides, 16:9 widescreen).

Theme: design philosophy + capability taxonomy + orchestration scenarios.
Reuses style helpers from build_pptx.py but is an entirely independent deck.

Run:  python3 build_spec_pptx.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# -------------------------------------------------------------------------
# Palette / constants  (kept identical to build_pptx.py for visual parity)
# -------------------------------------------------------------------------
NAVY = RGBColor(0x1F, 0x2A, 0x44)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
GRAY_BG = RGBColor(0xF3, 0xF4, 0xF6)
GRAY_MID = RGBColor(0xD1, 0xD5, 0xDB)
GRAY_TEXT = RGBColor(0x4B, 0x55, 0x63)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x11, 0x18, 0x27)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)

HEADER_TEXT = "Secmind · Atomic Investigation Skills"
DATE_STR = "2026-05-24"
FOOTER_TEXT = "secmind-investigator · Spec Design Review v1.0"
OUTPUT_PATH = Path(
    "/Users/asoiso/workspace/Secmind/docs/slides/"
    "2026-05-24-atomic-skills-spec.pptx"
)

FONT_BODY = "Microsoft YaHei"
FONT_MONO = "Menlo"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# -------------------------------------------------------------------------
# Low-level helpers (mirrors build_pptx.py)
# -------------------------------------------------------------------------
def add_rect(slide, left, top, width, height, fill, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if not line:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = GRAY_MID
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font=FONT_BODY):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_header_footer(slide, page_no, total):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.45), NAVY)
    add_text(slide, Inches(0.4), Inches(0.08), Inches(10), Inches(0.35),
             HEADER_TEXT, size=12, color=WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(10.5), Inches(0.08), Inches(2.6), Inches(0.35),
             DATE_STR, size=11, color=GRAY_MID,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(0), Inches(7.25), SLIDE_W, Inches(0.25), GRAY_BG)
    add_text(slide, Inches(0.4), Inches(7.25), Inches(8), Inches(0.25),
             FOOTER_TEXT, size=9, color=GRAY_TEXT,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(11.5), Inches(7.25), Inches(1.6), Inches(0.25),
             f"{page_no} / {total}", size=9, color=GRAY_TEXT,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_title(slide, title, subtitle=None):
    add_text(slide, Inches(0.5), Inches(0.65), Inches(12.3), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    add_rect(slide, Inches(0.5), Inches(1.35), Inches(0.9), Inches(0.06), BLUE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.45),
                 subtitle, size=14, color=GRAY_TEXT)


def add_bullets(slide, left, top, width, height, items, *,
                size=18, color=DARK_TEXT, line_spacing=1.25,
                bullet_color=BLUE, bullet="■ "):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        b = p.add_run()
        b.text = bullet
        b.font.name = FONT_BODY
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        r = p.add_run()
        r.text = item
        r.font.name = FONT_BODY
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_code_block(slide, left, top, width, height, code, *, size=13):
    add_rect(slide, left, top, width, height, RGBColor(0x0F, 0x17, 0x2A))
    tb = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), height - Inches(0.2),
    )
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = FONT_MONO
        run.font.size = Pt(size)
        run.font.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)


def add_table(slide, left, top, width, height, rows, *,
              header_fill=NAVY, header_text=WHITE, body_text=DARK_TEXT,
              font_size=12, header_size=13, col_widths=None):
    nrows = len(rows)
    ncols = len(rows[0])
    tbl_shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = tbl_shape.table
    if col_widths:
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = w
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                fc, fs, fb = header_text, header_size, True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 == 1 else GRAY_BG
                fc, fs, fb = body_text, font_size, False
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(0)
            tf.margin_right = Emu(0)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            if p.runs:
                run = p.runs[0]
                run.text = str(val)
            else:
                run = p.add_run()
                run.text = str(val)
            run.font.name = FONT_BODY
            run.font.size = Pt(fs)
            run.font.bold = fb
            run.font.color.rgb = fc
    return tbl


def add_callout(slide, left, top, width, height, text, *,
                fill=BLUE, color=WHITE, size=14, bold=True):
    add_rect(slide, left, top, width, height, fill)
    add_text(slide, left, top, width, height, text,
             size=size, bold=bold, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# -------------------------------------------------------------------------
# Slide builders
# -------------------------------------------------------------------------
def s01_cover(prs, total):
    s = new_slide(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    add_rect(s, Inches(0.5), Inches(1.9), Inches(0.18), Inches(3.7), BLUE)
    add_text(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.2),
             "安全告警调查 · 原子能力设计 (v1.0)",
             size=40, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(3.1), Inches(11.5), Inches(0.6),
             "Atomic Investigation Capabilities — Spec Design Review",
             size=22, color=GRAY_MID)
    add_text(s, Inches(0.9), Inches(3.85), Inches(11.5), Inches(0.5),
             "通用企业 SOC 调查动作的最小积木集",
             size=18, color=BLUE, bold=True)
    add_text(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(0.5),
             f"日期 · {DATE_STR}", size=16, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.4),
             "Secmind / secmind-investigator", size=13, color=GRAY_MID)
    add_text(s, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3),
             f"1 / {total}", size=10, color=GRAY_MID,
             align=PP_ALIGN.RIGHT)


def s02_agenda(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "议程", "围绕设计哲学 · 能力分类 · 编排场景三大主线")
    items = [
        ("01", "设计哲学 / 命名空间"),
        ("02", "抽象数据源清单"),
        ("03", "能力总览 7 大类"),
        ("04", "A-G 类原子能力精读"),
        ("05", "数据源 × 能力依赖矩阵"),
        ("06", "6 个典型告警编排场景"),
        ("07", "落地路线图 P0 / P1 / P2"),
        ("08", "总结 & Q&A"),
    ]
    box_w = Inches(2.95)
    box_h = Inches(2.1)
    gap_x = Inches(0.15)
    gap_y = Inches(0.3)
    start_x = Inches(0.5)
    start_y = Inches(2.3)
    for i, (idx, title) in enumerate(items):
        row = i // 4
        col = i % 4
        x = start_x + (box_w + gap_x) * col
        y = start_y + (box_h + gap_y) * row
        add_rect(s, x, y, box_w, box_h, GRAY_BG)
        add_rect(s, x, y, Inches(0.1), box_h, BLUE)
        add_text(s, x + Inches(0.3), y + Inches(0.25),
                 box_w - Inches(0.4), Inches(0.6),
                 idx, size=28, bold=True, color=BLUE)
        add_text(s, x + Inches(0.3), y + Inches(1.0),
                 box_w - Inches(0.4), Inches(1.0),
                 title, size=18, bold=True, color=NAVY)


def s03_philosophy(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "设计哲学:七条铁律", "Atomic / Composable / Decoupled / Observable / Structured / Privacy-explicit / Baseline-first")

    rules = [
        ("1", "原子化",        "单一职责,一个能力 = 一个调查动作",     BLUE),
        ("2", "可组合",        "任意串行 / 并行 / 嵌套,依赖显式声明",  BLUE),
        ("3", "数据源解耦",    "按语义命名,绑定抽象 DS-XXX 而非具体厂商", GREEN),
        ("4", "降级可观测",    "数据源不可达必须 partial=true + 显式缺口", AMBER),
        ("5", "输出结构化",    "JSON-shaped,可被另一个能力直接消费",     PURPLE),
        ("6", "隐私显式",      "声明 PII 字段,上层按权限裁剪",         RED),
        ("7", "基线优先于规则", "判断'异常'优先于'匹配 IOC'",        NAVY),
    ]
    box_w = Inches(6.05)
    box_h = Inches(0.65)
    gap_y = Inches(0.1)
    top = Inches(2.0)
    for i, (num, name, desc, col) in enumerate(rules):
        row = i // 2
        c = i % 2
        x = Inches(0.5) + (box_w + Inches(0.25)) * c
        y = top + (box_h + gap_y) * row
        add_rect(s, x, y, box_w, box_h, GRAY_BG)
        add_rect(s, x, y, Inches(0.55), box_h, col)
        add_text(s, x, y, Inches(0.55), box_h, num,
                 size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.7), y + Inches(0.08),
                 Inches(1.6), Inches(0.5),
                 name, size=14, bold=True, color=col,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(2.3), y + Inches(0.08),
                 box_w - Inches(2.4), Inches(0.5),
                 desc, size=12, color=DARK_TEXT,
                 anchor=MSO_ANCHOR.MIDDLE)

    add_callout(s, Inches(0.5), Inches(6.55), Inches(12.35), Inches(0.55),
                "原子化是地基,LLM Agent 是大脑",
                fill=NAVY, color=WHITE, size=15, bold=True)


def s04_namespace(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "命名空间体系",
              "investigate.<category>.<entity|activity>.<action>")

    # format strip
    add_rect(s, Inches(0.5), Inches(2.0), Inches(12.35), Inches(0.65), NAVY)
    add_text(s, Inches(0.5), Inches(2.0), Inches(12.35), Inches(0.65),
             "investigate.<category>.<entity|activity>.<action>",
             size=18, bold=True, color=WHITE, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    rows = [
        ["缩写", "Category", "含义", "示例 namespace"],
        ["A", "entity",   "实体画像",      "investigate.entity.user.profile"],
        ["B", "activity", "行为活动",      "investigate.activity.login"],
        ["C", "relation", "关联推理",      "investigate.relation.lateral_movement"],
        ["D", "context",  "上下文与基线",  "investigate.context.baseline.self"],
        ["E", "intel",    "外部情报",      "investigate.intel.ip.reputation"],
        ["F", "verdict",  "判定与产出",    "investigate.verdict.tp_fp"],
        ["G", "meta",     "元能力",        "investigate.meta.next_step"],
    ]
    add_table(s, Inches(0.5), Inches(2.85), Inches(12.35), Inches(4.0), rows,
              font_size=12, header_size=13,
              col_widths=[Inches(1.0), Inches(1.8), Inches(2.4), Inches(7.15)])


def s05_capability_card(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "能力卡统一模板",
              "每张能力卡 = 一份契约,可被 LLM 直接读懂")

    add_bullets(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(4.6),
                [
                    "ID(A1 / B3 ...)",
                    "名称 / 命名空间",
                    "目的(一句话,LLM 友好)",
                    "输入字段表",
                    "去哪里查 → 查什么",
                    "输出 JSON 字段",
                    "典型编排触发场景",
                    "依赖能力 / 失败降级 / 备注",
                ], size=14, line_spacing=1.45)

    code = (
        "id: A1\n"
        "name: 用户/账户画像\n"
        "namespace: investigate.entity.user.profile\n"
        "purpose: 给定一个 principal,产出该用户的多源画像\n"
        "inputs:\n"
        "  principal: str\n"
        "  time_window: TimeWindow\n"
        "sources:\n"
        "  - DS-HR    : department / leaver_in_30d\n"
        "  - DS-IDP   : status / mfa / privileged_groups\n"
        "  - DS-CMDB  : devices / saas_apps\n"
        "outputs:\n"
        "  identity / auth / privilege / assets / email\n"
        "  risk_tags / prior_alerts_30d\n"
        "on_partial:\n"
        "  partial=true · partial_reasons=['DS-HR unavailable']\n"
        "pii: [identity.email, hr.manager]"
    )
    add_code_block(s, Inches(6.5), Inches(2.0), Inches(6.35), Inches(4.8),
                   code, size=12)

    add_callout(s, Inches(0.5), Inches(6.85), Inches(12.35), Inches(0.35),
                "结构统一 → LLM 一次学习,所有 83 个能力即插即用",
                fill=GREEN, color=WHITE, size=12)


def s06_ds_p0(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "抽象数据源清单 · P0 核心 9 个",
              "按语义命名,落地时各企业做一次 DS-X → 厂商接口的映射")

    rows = [
        ["抽象 DS", "含义", "典型实际系统", "保留期"],
        ["DS-IDP",   "身份与认证",       "Entra ID / Okta / AD",              "90-365 d"],
        ["DS-HR",    "人事系统",         "Workday / SAP / 飞书人事",          "permanent"],
        ["DS-CMDB",  "资产配置",         "ServiceNow / Lansweeper",           "permanent"],
        ["DS-EDR",   "终端检测响应",     "CrowdStrike / SentinelOne / 360",   "30-180 d"],
        ["DS-SIEM",  "日志聚合与告警",   "Splunk / Sentinel / Chronicle",     "90-365 d"],
        ["DS-EMAIL", "邮件安全",         "M365 / Proofpoint / 网易企业邮",    "90-365 d"],
        ["DS-FW",    "网络防火墙",       "Palo Alto / Fortinet / 华为",       "30-180 d"],
        ["DS-TI",    "威胁情报",         "VirusTotal / OTX / 商业 TI",        "API only"],
        ["DS-KB",    "知识库 / Runbook", "Confluence / Notion / 内部 Wiki",   "permanent"],
    ]
    add_table(s, Inches(0.5), Inches(2.0), Inches(12.35), Inches(4.6), rows,
              font_size=12, header_size=13,
              col_widths=[Inches(1.5), Inches(2.5), Inches(5.85), Inches(2.5)])

    add_callout(s, Inches(0.5), Inches(6.75), Inches(12.35), Inches(0.4),
                "落地时各企业做一次 DS-X → 厂商接口的映射,能力代码完全不动",
                fill=NAVY, color=WHITE, size=12)


def s07_ds_p1_p2(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "抽象数据源清单 · P1 + P2 扩展",
              "P1 应对云/横移/外泄,P2 覆盖归因/合规/供应链")

    rows_p1 = [
        ["P1 (9 类)", "含义"],
        ["DS-NDR",     "网络检测响应"],
        ["DS-CLOUD",   "云审计日志"],
        ["DS-DLP",     "数据防泄漏"],
        ["DS-CASB",    "SaaS 治理"],
        ["DS-PAM",     "特权访问管理"],
        ["DS-DAM",     "数据库访问监控"],
        ["DS-SANDBOX", "沙箱动态分析"],
        ["DS-VULN",    "漏洞扫描"],
        ["DS-WHOIS",   "域名注册信息"],
    ]
    add_text(s, Inches(0.5), Inches(1.95), Inches(6.0), Inches(0.4),
             "P1 · 场景扩展", size=14, bold=True, color=GREEN)
    add_table(s, Inches(0.5), Inches(2.4), Inches(6.0), Inches(4.4), rows_p1,
              font_size=11, header_size=12,
              col_widths=[Inches(2.0), Inches(4.0)])

    rows_p2 = [
        ["P2 (10 类)", "含义"],
        ["DS-CODE",    "代码仓库与构建"],
        ["DS-MOBILE",  "移动设备管理"],
        ["DS-DNS",     "递归 DNS 与 pDNS"],
        ["DS-FILE",    "文件存储审计"],
        ["DS-WAF",     "应用层防护"],
        ["DS-VPN",     "远程接入"],
        ["DS-PROXY",   "上网代理"],
        ["DS-TICKET",  "工单变更"],
        ["DS-AV",      "防病毒事件"],
        ["DS-IPAM",    "IP 资产管理"],
    ]
    add_text(s, Inches(6.9), Inches(1.95), Inches(6.0), Inches(0.4),
             "P2 · 深度专题", size=14, bold=True, color=PURPLE)
    add_table(s, Inches(6.9), Inches(2.4), Inches(6.0), Inches(4.6), rows_p2,
              font_size=11, header_size=12,
              col_widths=[Inches(2.0), Inches(4.0)])


def s08_overview(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "能力总览 · 7 大类 · 83 个原子能力",
              "类目化设计,LLM 按需挑选")

    cats = [
        ("A", "实体画像",       14, BLUE),
        ("B", "行为活动",       17, GREEN),
        ("C", "关联推理",       10, AMBER),
        ("D", "上下文与基线",   11, PURPLE),
        ("E", "外部情报",       12, RED),
        ("F", "判定与产出",     11, NAVY),
        ("G", "元能力",          8, GRAY_TEXT),
    ]
    # column bars
    x0 = Inches(0.5)
    w_total = Inches(12.35)
    n = len(cats)
    col_w = Inches((13.333 - 1.0) / n - 0.15)
    gap = Inches(0.15)
    top = Inches(2.2)
    chart_h = Inches(3.5)
    max_count = max(c[2] for c in cats)
    for i, (k, name, count, col) in enumerate(cats):
        x = x0 + (col_w + gap) * i
        bar_h_in = 3.5 * count / max_count
        bar_top = top + Inches(3.5 - bar_h_in)
        # background slot
        add_rect(s, x, top, col_w, chart_h, GRAY_BG)
        # bar
        add_rect(s, x, bar_top, col_w, Inches(bar_h_in), col)
        # count on bar
        add_text(s, x, bar_top + Inches(0.05), col_w, Inches(0.5),
                 str(count), size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        # label
        add_text(s, x, top + chart_h + Inches(0.1), col_w, Inches(0.4),
                 f"{k} · {name}", size=12, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_callout(s, Inches(0.5), Inches(6.6), Inches(12.35), Inches(0.55),
                "总计 83 个原子能力 · 覆盖从'拿数'到'判定'到'下一步'的完整调查链",
                fill=NAVY, color=WHITE, size=14, bold=True)


def s09_a1_card(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "A 类 · 实体画像 · A1 精读",
              "investigate.entity.user.profile")

    rows = [
        ["项", "值"],
        ["输入", "principal: str  ·  time_window: TimeWindow"],
        ["数据源", "DS-HR● DS-IDP● DS-CMDB○ DS-PAM○ DS-EMAIL○ DS-KB○ 观察名单○"],
        ["输出字段", "identity / auth / privilege / assets / email / risk_tags / prior_alerts_30d"],
        ["失败降级", "HR 不可达 → partial=true · partial_reasons=['DS-HR unavailable']"],
    ]
    add_table(s, Inches(0.5), Inches(2.0), Inches(12.35), Inches(2.4), rows,
              font_size=12, header_size=13,
              col_widths=[Inches(1.4), Inches(10.95)])

    # risk tags showcase
    add_text(s, Inches(0.5), Inches(4.55), Inches(12.35), Inches(0.4),
             "典型 risk_tags 示例", size=14, bold=True, color=NAVY)
    tags = [
        ("leaver_30d",      RED),
        ("watchlist",       RED),
        ("dormant_active",  AMBER),
        ("priv_no_mfa",     AMBER),
        ("contractor",      BLUE),
        ("vip",             PURPLE),
        ("recent_password_reset", GRAY_TEXT),
    ]
    x = Inches(0.5)
    y = Inches(5.0)
    for label, col in tags:
        w = Inches(0.25 + 0.13 * len(label))
        add_rect(s, x, y, w, Inches(0.45), col)
        add_text(s, x, y, w, Inches(0.45), label,
                 size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x = x + w + Inches(0.12)

    add_callout(s, Inches(0.5), Inches(5.85), Inches(12.35), Inches(1.05),
                "备注:HR 数据 T+1 滞后 · 邮箱转发规则是 BEC 持久化经典点\n"
                "LLM 直接据此判断:是否高敏感用户 / 是否启用 MFA / 是否离职态",
                fill=NAVY, color=WHITE, size=13, bold=True)


def s10_a_others(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "A 类 · 其他 13 个实体能力一览",
              "实体是调查的锚点,所有 B/C 能力都围绕实体展开")

    items = [
        ("A2",  "主机 / 终端",       BLUE),
        ("A3",  "服务器 / 资产",     BLUE),
        ("A4",  "IP 地址",           BLUE),
        ("A5",  "域名",              BLUE),
        ("A6",  "URL",               BLUE),
        ("A7",  "文件 / Hash",       BLUE),
        ("A8",  "进程",              BLUE),
        ("A9",  "邮件地址",          BLUE),
        ("A10", "证书 / 公钥",       GREEN),
        ("A11", "云资源",            GREEN),
        ("A12", "服务账户 / API Key", GREEN),
        ("A13", "移动设备",          GREEN),
        ("A14", "DB 账户 / 对象",    GREEN),
    ]
    # 4 columns x ceil(13/4)=4 rows
    box_w = Inches(2.95)
    box_h = Inches(0.95)
    gap_x = Inches(0.15)
    gap_y = Inches(0.15)
    x0 = Inches(0.5)
    y0 = Inches(2.05)
    for i, (idx, name, col) in enumerate(items):
        r = i // 4
        c = i % 4
        x = x0 + (box_w + gap_x) * c
        y = y0 + (box_h + gap_y) * r
        add_rect(s, x, y, box_w, box_h, GRAY_BG)
        add_rect(s, x, y, Inches(0.6), box_h, col)
        add_text(s, x, y, Inches(0.6), box_h, idx,
                 size=16, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.75), y, box_w - Inches(0.85), box_h,
                 name, size=14, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)

    add_callout(s, Inches(0.5), Inches(6.7), Inches(12.35), Inches(0.45),
                "命名一致性:所有 A 类输出 entity_id / entity_type / risk_tags / prior_alerts_30d",
                fill=NAVY, color=WHITE, size=12)


def s11_b_category(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "B 类 · 行为活动 · 17 个能力",
              "按数据源/语义分四组")

    groups = [
        ("终端类",
         "B3 进程执行 · B6 文件操作 · B7 配置变更 · "
         "B15 持久化 · B16 凭据访问 · B17 防御规避",
         BLUE),
        ("网络类",
         "B4 网络连接 · B5 DNS · B11 USB 接入 · B12 浏览器活动",
         GREEN),
        ("身份 / 通信类",
         "B1 登录 · B2 权限变更 · B8 邮件收发 · B13 SaaS OAuth 授权",
         AMBER),
        ("数据 / 云 / 容器类",
         "B9 云控制面 · B10 数据访问 · B14 容器 / K8s 事件",
         PURPLE),
    ]
    top = Inches(2.0)
    h = Inches(1.05)
    gap = Inches(0.12)
    for i, (g, items, col) in enumerate(groups):
        y = top + (h + gap) * i
        add_rect(s, Inches(0.5), y, Inches(12.35), h, GRAY_BG)
        add_rect(s, Inches(0.5), y, Inches(0.18), h, col)
        add_text(s, Inches(0.85), y + Inches(0.1),
                 Inches(3.0), Inches(0.4),
                 g, size=18, bold=True, color=col)
        add_text(s, Inches(0.85), y + Inches(0.5),
                 Inches(11.4), Inches(0.5),
                 items, size=13, color=DARK_TEXT)

    add_callout(s, Inches(0.5), Inches(6.7), Inches(12.35), Inches(0.45),
                "B 类回答:这个实体在窗口期内做了什么 / 被做了什么",
                fill=NAVY, color=WHITE, size=12)


def s12_c_category(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "C 类 · 关联推理 · 10 个能力",
              "把孤立证据拼成攻击叙事")

    rows = [
        ["ID", "能力", "一句话作用"],
        ["C1",  "时间窗共现",    "同一实体在 ±Δt 内的多源事件聚合"],
        ["C2",  "实体关系图",    "user / host / ip / domain 之间的边"],
        ["C3",  "父子进程链",    "完整 process lineage"],
        ["C4",  "ATT&CK 阶段定位",  "事件 → tactic / technique 映射"],
        ["C5",  "横向移动追踪",  "登录链 + 远控工具调用链"],
        ["C6",  "扇出 / 扇入",   "1→N / N→1 的爆破 / 内传模式"],
        ["C7",  "同源事件",      "相同 IOC / 同一 campaign 的事件聚类"],
        ["C8",  "跨源拼接",      "EDR + 邮件 + 防火墙的同一事件还原"],
        ["C9",  "受害者扩展",    "同邮件模板 / 同样本 / 同 C2 的更多受害者"],
        ["C10", "异常路径检测",  "认证 / 进程 / 网络的稀有路径"],
    ]
    add_table(s, Inches(0.5), Inches(2.0), Inches(12.35), Inches(4.6), rows,
              font_size=11, header_size=12,
              col_widths=[Inches(1.0), Inches(3.0), Inches(8.35)])

    add_callout(s, Inches(0.5), Inches(6.75), Inches(12.35), Inches(0.4),
                "关联类把孤立证据拼成攻击叙事 —— Verdict 的真正前提",
                fill=NAVY, color=WHITE, size=12)


def s13_d_category(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "D 类 · 上下文与基线 · 11 个能力",
              "基线优先于规则 —— 异常检测的真实地基")

    # group 1 - business/identity
    add_rect(s, Inches(0.5), Inches(2.0), Inches(4.0), Inches(4.6), GRAY_BG)
    add_rect(s, Inches(0.5), Inches(2.0), Inches(4.0), Inches(0.4), BLUE)
    add_text(s, Inches(0.5), Inches(2.0), Inches(4.0), Inches(0.4),
             "业务 / 身份",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(0.65), Inches(2.55), Inches(3.7), Inches(4.0),
                ["D1 资产业务上下文",
                 "D2 用户身份上下文",
                 "D3 漏洞 / 暴露上下文",
                 "D4 合规 / 合同上下文"],
                size=12, line_spacing=1.5)

    # group 2 - baseline
    add_rect(s, Inches(4.7), Inches(2.0), Inches(4.0), Inches(4.6), GRAY_BG)
    add_rect(s, Inches(4.7), Inches(2.0), Inches(4.0), Inches(0.4), GREEN)
    add_text(s, Inches(4.7), Inches(2.0), Inches(4.0), Inches(0.4),
             "基线对比 (核心)",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(4.85), Inches(2.55), Inches(3.7), Inches(4.0),
                ["D5 自身历史基线",
                 "D6 同组 / 同岗对等基线",
                 "D7 全组织罕见度",
                 "(回答:对这个人/这台机, 是不是异常?)"],
                size=12, line_spacing=1.5)

    # group 3 - env
    add_rect(s, Inches(8.9), Inches(2.0), Inches(3.95), Inches(4.6), GRAY_BG)
    add_rect(s, Inches(8.9), Inches(2.0), Inches(3.95), Inches(0.4), AMBER)
    add_text(s, Inches(8.9), Inches(2.0), Inches(3.95), Inches(0.4),
             "环境维度",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(9.05), Inches(2.55), Inches(3.7), Inches(4.0),
                ["D8 业务时间窗",
                 "D9 地理位置基线",
                 "D10 设备指纹",
                 "D11 网络位置"],
                size=12, line_spacing=1.5)

    add_callout(s, Inches(0.5), Inches(6.75), Inches(12.35), Inches(0.4),
                "强调:基线优先于规则 —— 异常检测的真实地基",
                fill=NAVY, color=WHITE, size=13, bold=True)


def s14_e_category(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "E 类 · 外部情报 · 12 个能力",
              "情报是外援,基线是底座")

    groups = [
        ("信誉类",       "E1 IP · E2 域名/URL · E3 Hash",                      BLUE),
        ("历史 / 规则",  "E4 WHOIS + pDNS + CT · E6 YARA / Sigma",             GREEN),
        ("动态分析",     "E5 沙箱引爆 (静/动态结合)",                          AMBER),
        ("映射 / 归因",  "E7 TTP → MITRE · E8 CVE · E9 APT 归因",              PURPLE),
        ("泄露 / 品牌 / 供应链", "E10 凭据泄露 · E11 仿冒情报 · E12 供应链情报",  RED),
    ]
    top = Inches(2.0)
    h = Inches(0.85)
    gap = Inches(0.1)
    for i, (g, items, col) in enumerate(groups):
        y = top + (h + gap) * i
        add_rect(s, Inches(0.5), y, Inches(12.35), h, GRAY_BG)
        add_rect(s, Inches(0.5), y, Inches(0.18), h, col)
        add_text(s, Inches(0.85), y + Inches(0.1),
                 Inches(3.0), Inches(0.35),
                 g, size=16, bold=True, color=col)
        add_text(s, Inches(0.85), y + Inches(0.45),
                 Inches(11.4), Inches(0.4),
                 items, size=13, color=DARK_TEXT)

    add_callout(s, Inches(0.5), Inches(6.75), Inches(12.35), Inches(0.4),
                "原则:E 类只补强 A/B/C/D 的结论,不单独充当 Verdict",
                fill=NAVY, color=WHITE, size=12)


def s15_f_category(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "F 类 · 判定与产出 · 11 个能力",
              "判定是终点 · 改进是回路")

    groups = [
        ("判定",     "F1 TP / FP · F2 影响范围 · F3 攻击阶段总结",                BLUE),
        ("提取",     "F4 IOC 提取与标准化 · F5 攻击者画像",                       GREEN),
        ("重构",     "F6 时间线 · F7 受影响实体清单 · F9 根因 / 初始入口",        PURPLE),
        ("输出",     "F8 处置建议 · F10 通报材料 · F11 检测改进建议",             AMBER),
    ]
    top = Inches(2.0)
    h = Inches(1.0)
    gap = Inches(0.12)
    for i, (g, items, col) in enumerate(groups):
        y = top + (h + gap) * i
        add_rect(s, Inches(0.5), y, Inches(12.35), h, GRAY_BG)
        add_rect(s, Inches(0.5), y, Inches(0.18), h, col)
        add_text(s, Inches(0.85), y + Inches(0.1),
                 Inches(3.0), Inches(0.4),
                 g, size=18, bold=True, color=col)
        add_text(s, Inches(0.85), y + Inches(0.5),
                 Inches(11.4), Inches(0.5),
                 items, size=13, color=DARK_TEXT)

    add_callout(s, Inches(0.5), Inches(6.7), Inches(12.35), Inches(0.45),
                "F11 把每次调查的洞察反哺成新的检测规则 —— 闭环改进",
                fill=NAVY, color=WHITE, size=13, bold=True)


def s16_g_category(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "G 类 · 元能力 · 8 个能力",
              "元能力让 Agent 学会决定'下一步查什么'")

    items = [
        ("G1", "假设生成",                 BLUE),
        ("G2", "反证维护",                 BLUE),
        ("G3", "不确定性 + 下一步推荐",    GREEN),
        ("G4", "历史相似告警回查",         GREEN),
        ("G5", "调查路径优化",             AMBER),
        ("G6", "证据链完整性自检",         AMBER),
        ("G7", "PII / 合规边界检查",       RED),
        ("G8", "人在环路升级",             PURPLE),
    ]
    box_w = Inches(2.95)
    box_h = Inches(1.05)
    gap_x = Inches(0.15)
    gap_y = Inches(0.2)
    x0 = Inches(0.5)
    y0 = Inches(2.05)
    for i, (idx, name, col) in enumerate(items):
        r = i // 4
        c = i % 4
        x = x0 + (box_w + gap_x) * c
        y = y0 + (box_h + gap_y) * r
        add_rect(s, x, y, box_w, box_h, GRAY_BG)
        add_rect(s, x, y, Inches(0.7), box_h, col)
        add_text(s, x, y, Inches(0.7), box_h, idx,
                 size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.85), y, box_w - Inches(0.95), box_h,
                 name, size=14, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)

    add_callout(s, Inches(0.5), Inches(6.55), Inches(12.35), Inches(0.55),
                "G 类 = Agent 自我反思 · 隐私守门 · 升级触发",
                fill=NAVY, color=WHITE, size=14, bold=True)


def s17_matrix(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "数据源 × 能力依赖矩阵(节选)",
              "● 必需  ○ 可选  ·  完整矩阵见 spec §3")

    rows = [
        ["能力 \\ DS",          "IDP", "HR", "CMDB", "EDR", "SIEM", "EMAIL", "FW", "TI", "KB"],
        ["A1 用户画像",          "●",   "●",  "○",    "",    "",     "○",    "",    "",   "○"],
        ["A2 主机画像",          "",    "",   "●",    "●",   "○",    "",     "",    "",   "○"],
        ["A4 IP 画像",           "",    "",   "○",    "",    "●",    "",     "●",   "○",  ""],
        ["A7 文件 Hash",         "",    "",   "",     "●",   "",     "○",    "",    "●",  ""],
        ["B1 登录活动",          "●",   "",   "",     "",    "●",    "",     "",    "",   ""],
        ["B3 进程执行",          "",    "",   "",     "●",   "○",    "",     "",    "",   ""],
        ["B4 网络连接",          "",    "",   "○",    "○",   "○",    "",     "●",   "",   ""],
        ["B8 邮件收发",          "",    "",   "",     "",    "",     "●",    "",    "",   ""],
        ["E1 IP 信誉",           "",    "",   "",     "",    "",     "",     "",    "●",  ""],
        ["F8 处置建议",          "○",   "",   "○",    "○",   "○",    "○",    "○",   "○",  "●"],
    ]
    add_table(s, Inches(0.5), Inches(2.0), Inches(12.35), Inches(4.7), rows,
              font_size=10, header_size=11,
              col_widths=[Inches(2.35)] + [Inches(1.11)] * 9)

    add_callout(s, Inches(0.5), Inches(6.8), Inches(12.35), Inches(0.4),
                "矩阵驱动 P0 数据源选型:9 个 DS 即可支撑前 25 个能力",
                fill=BLUE, color=WHITE, size=12)


def _scenario_flow(slide, top, steps, accent=BLUE):
    """Render a numbered vertical flow inside a column."""
    y = top
    for i, text in enumerate(steps, 1):
        add_rect(slide, Inches(0.55), y, Inches(0.42), Inches(0.36), accent)
        add_text(slide, Inches(0.55), y, Inches(0.42), Inches(0.36), str(i),
                 size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(1.05), y - Inches(0.02),
                 Inches(11.8), Inches(0.42),
                 text, size=12, color=DARK_TEXT,
                 anchor=MSO_ANCHOR.MIDDLE)
        y = y + Inches(0.42)


def s18_phishing(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "编排示例 1 · 钓鱼邮件告警",
              "入口:EMAIL 网关告警 “外部邮件含可疑链接”")
    steps = [
        "A9 邮件地址画像 + A1 用户画像 (锚点)",
        "B8 邮件收发 (完整 header / 附件 / 链接)",
        "A6 URL 画像 + E2 URL/Domain 信誉 → 未知则 E5 沙箱引爆",
        "A7 附件 Hash + E3 Hash 信誉 → 未知则 E5 沙箱",
        "D9 地理基线 + E11 仿冒情报 (是否伪装合作伙伴)",
        "C9 受害者扩展 (同模板 / 同发件人)",
        "B12 浏览器活动 (收件人是否实际点击)",
        "C1 时间窗共现 (±2h 内同实体多源事件)",
        "F1 判定 → F4 IOC 提取 → F8 处置建议 (退/隔/封)",
        "G4 历史相似回查 (同样本/同 campaign 历史告警)",
    ]
    _scenario_flow(s, Inches(2.0), steps, accent=BLUE)


def s19_ransomware(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "编排示例 2 · 勒索软件告警",
              "入口:EDR “可疑加密行为”")
    steps = [
        "A2 主机画像 + A1 用户画像 (受害者锚点)",
        "A8 进程画像 + A7 镜像 Hash 画像",
        "B3 进程执行 (±10min 进程树还原)",
        "B6 文件操作 (大量加密改写 / 勒索信落地)",
        "B17 防御规避 (关闭 AV/EDR · 清除卷影副本)",
        "B15 持久化 (新增服务 / 计划任务 / 启动项)",
        "C5 横向移动 + C9 受害者扩展 (同样本同 C2 的其他主机)",
        "D1 资产业务上下文 (业务等级 / 备份策略)",
        "F4 IOC 标准化 → F8 处置建议 (隔离 / 重置 / 通报)",
        "G8 人在环路升级 (必升级 · 业务影响重大)",
    ]
    _scenario_flow(s, Inches(2.0), steps, accent=RED)


def s20_exfil_mining(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "编排示例 3 + 4 · 数据外传 / 挖矿",
              "双场景对比 · 体现基线优先 + 多渠道关联")

    # left col - data exfil
    add_rect(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(4.7), GRAY_BG)
    add_rect(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(0.45), AMBER)
    add_text(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(0.45),
             "数据外传 (Insider / BEC / 外部入侵)",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(0.65), Inches(2.6), Inches(5.7), Inches(4.0),
                [
                    "A1 + D2 离职倒计时 / 高敏感岗位",
                    "B10 数据访问 (DLP / DAM / CASB)",
                    "B8 + B12 + B11 + B13 全渠道外传",
                    "A4 + A6 流向 → 外部 IP / SaaS",
                    "D5 + D6 历史与同组基线",
                    "F1 判定三选一:Insider / 误操 / 入侵",
                    "F8 处置 · G7 隐私守门",
                ],
                size=12, line_spacing=1.45)

    # right col - mining
    add_rect(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(4.7), GRAY_BG)
    add_rect(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(0.45), GREEN)
    add_text(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(0.45),
             "挖矿木马",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(7.0), Inches(2.6), Inches(5.7), Inches(4.0),
                [
                    "A4 IP 画像 + E1 信誉 (矿池标签)",
                    "A2 主机 + A8 进程 (xmrig / 父进程)",
                    "B3 进程执行 (CPU 占用 / 持久化)",
                    "B4 网络连接 (Stratum / 矿池端口)",
                    "D7 全组织罕见度 (此进程是否仅此一例)",
                    "C9 受害者扩展 (同矿池的其他主机)",
                    "F9 根因 (初始入口) → F8 处置",
                ],
                size=12, line_spacing=1.45)

    add_callout(s, Inches(0.5), Inches(6.8), Inches(12.35), Inches(0.4),
                "共性:基线 (D5/D6/D7) 是从'有告警'到'真异常'的关键一步",
                fill=NAVY, color=WHITE, size=12)


def s21_bruteforce_lateral(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "编排示例 5 + 6 · 暴破 / 凭据填充 + 横向移动",
              "高频实战场景 · 体现 C 类关联 + G 类升级")

    # left col - bruteforce
    add_rect(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(4.7), GRAY_BG)
    add_rect(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(0.45), BLUE)
    add_text(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(0.45),
             "暴破 / 凭据填充",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(0.65), Inches(2.6), Inches(5.7), Inches(4.0),
                [
                    "A1 + A4 (账户 + 源 IP 锚点)",
                    "B1 登录 (成功/失败 分布 / UA / Geo)",
                    "D9 + D11 地理 / 网络位置基线",
                    "E10 凭据泄露 (是否命中泄露库)",
                    "C6 扇出:单 IP→多账户 / 多 IP→单账户",
                    "C1 时间窗 + D2 高权限 / VIP 加权",
                    "F8:临时封 IP + 强制重置 + 强制 MFA",
                ],
                size=12, line_spacing=1.45)

    # right col - lateral
    add_rect(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(4.7), GRAY_BG)
    add_rect(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(0.45), PURPLE)
    add_text(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(0.45),
             "横向移动",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(7.0), Inches(2.6), Inches(5.7), Inches(4.0),
                [
                    "A2 + A1 (主机 + 账户)",
                    "B16 凭据访问 (LSASS / Mimikatz)",
                    "B3 进程执行 (psexec / wmic / winrm)",
                    "C3 父子链 + C5 完整横移路径",
                    "C2 已触达的全部资产图",
                    "C4 + E9 阶段 / APT 归因",
                    "F2 影响评估 → F8 (批量隔离 / 重置 KRBTGT)",
                    "G8 必升级",
                ],
                size=12, line_spacing=1.45)

    add_callout(s, Inches(0.5), Inches(6.8), Inches(12.35), Inches(0.4),
                "横向移动 = C 类关联能力的硬核试金石",
                fill=NAVY, color=WHITE, size=12)


def s22_p0(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "落地路线图 · P0 (MVP)",
              "闭环三类告警 · 25 个能力 · 9 个数据源")

    # left: targets + abilities
    add_text(s, Inches(0.5), Inches(2.0), Inches(8.5), Inches(0.4),
             "P0 目标场景", size=14, bold=True, color=NAVY)
    add_bullets(s, Inches(0.5), Inches(2.45), Inches(8.5), Inches(1.4),
                ["钓鱼邮件 / 恶意软件 / 异常登录 —— 3 类告警端到端可闭环"],
                size=13, line_spacing=1.4)

    add_text(s, Inches(0.5), Inches(3.65), Inches(8.5), Inches(0.4),
             "25 个能力清单", size=14, bold=True, color=NAVY)
    rows = [
        ["类别", "能力"],
        ["A 实体 (5)",       "A1 A2 A4 A5 A7"],
        ["B 活动 (4)",       "B1 B3 B4 B8"],
        ["C 关联 (3)",       "C1 C2 C3"],
        ["D 上下文 (3)",     "D1 D2 D5"],
        ["E 情报 (4)",       "E1 E2 E3 E7"],
        ["F 判定 (4)",       "F1 F2 F4 F8"],
        ["G 元能力 (2)",     "G3 G4"],
    ]
    add_table(s, Inches(0.5), Inches(4.1), Inches(8.5), Inches(2.7), rows,
              font_size=11, header_size=12,
              col_widths=[Inches(2.0), Inches(6.5)])

    # right: 9 DS
    add_text(s, Inches(9.3), Inches(2.0), Inches(3.55), Inches(0.4),
             "9 个数据源", size=14, bold=True, color=GREEN)
    ds = ["DS-IDP", "DS-HR", "DS-CMDB", "DS-EDR", "DS-SIEM",
          "DS-EMAIL", "DS-FW", "DS-TI", "DS-KB"]
    for i, d in enumerate(ds):
        y = Inches(2.5) + Inches(0.5) * i
        add_rect(s, Inches(9.3), y, Inches(3.55), Inches(0.4), GRAY_BG)
        add_rect(s, Inches(9.3), y, Inches(0.1), Inches(0.4), GREEN)
        add_text(s, Inches(9.5), y, Inches(3.35), Inches(0.4),
                 d, size=12, bold=True, color=NAVY,
                 font=FONT_MONO, anchor=MSO_ANCHOR.MIDDLE)

    add_callout(s, Inches(0.5), Inches(6.9), Inches(12.35), Inches(0.3),
                "MVP 验收:跑通钓鱼 + 勒索 + 异常登录的三条端到端",
                fill=BLUE, color=WHITE, size=11)


def s23_p1_p2(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "路线图 · P1 + P2",
              "P1 场景扩展(+27) · P2 深度专题(+31)")

    # P1
    add_rect(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(4.4), GRAY_BG)
    add_rect(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(0.45), GREEN)
    add_text(s, Inches(0.5), Inches(2.0), Inches(6.0), Inches(0.45),
             "P1 · 云 / 外泄 / 勒索 / 横向",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(0.65), Inches(2.6), Inches(5.7), Inches(3.7),
                [
                    "新增 27 个能力",
                    "新增 DS-CLOUD / DLP / CASB / NDR",
                    "新增 DS-DAM / PAM / SANDBOX",
                    "新增 DS-VULN / WHOIS",
                    "聚焦:云审计 / 横移 / 数据外泄 / 特权滥用",
                ],
                size=13, line_spacing=1.5)

    # P2
    add_rect(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(4.4), GRAY_BG)
    add_rect(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(0.45), PURPLE)
    add_text(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(0.45),
             "P2 · 归因 / 供应链 / 内幕 / 合规",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(s, Inches(7.0), Inches(2.6), Inches(5.7), Inches(3.7),
                [
                    "新增 31 个能力",
                    "新增 DS-CODE / MOBILE / DNS / FILE",
                    "新增 DS-WAF / VPN / PROXY",
                    "新增 DS-TICKET / AV / IPAM",
                    "聚焦:APT 归因 / 供应链 / 内幕威胁 / 合规",
                ],
                size=13, line_spacing=1.5)

    # dependency timeline
    add_rect(s, Inches(0.5), Inches(6.55), Inches(12.35), Inches(0.65), NAVY)
    add_text(s, Inches(0.5), Inches(6.55), Inches(12.35), Inches(0.65),
             "DS 接入序:CMDB+HR+IDP → EDR+SIEM → EMAIL+FW → "
             "TI+KB → NDR+CLOUD+DLP → SANDBOX+WHOIS",
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def s24_summary(prs, page, total):
    s = new_slide(prs)
    add_header_footer(s, page, total)
    add_title(s, "总结 & Q&A",
              "Atomic Investigation Capabilities · Spec v1.0")

    # KPI tiles
    tiles = [
        ("83",  "原子能力",        BLUE),
        ("28",  "抽象数据源类目",  GREEN),
        ("7",   "能力大类",         AMBER),
        ("3",   "落地阶段 P0-P2",   PURPLE),
    ]
    tile_w = Inches(2.95)
    tile_h = Inches(1.6)
    gap = Inches(0.15)
    x0 = Inches(0.5)
    y0 = Inches(2.0)
    for i, (n, lbl, col) in enumerate(tiles):
        x = x0 + (tile_w + gap) * i
        add_rect(s, x, y0, tile_w, tile_h, GRAY_BG)
        add_rect(s, x, y0, tile_w, Inches(0.1), col)
        add_text(s, x, y0 + Inches(0.2), tile_w, Inches(0.8),
                 n, size=44, bold=True, color=col,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x, y0 + Inches(1.05), tile_w, Inches(0.4),
                 lbl, size=14, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # principles
    add_text(s, Inches(0.5), Inches(3.85), Inches(12.35), Inches(0.4),
             "关键原则", size=16, bold=True, color=NAVY)
    add_bullets(s, Inches(0.5), Inches(4.3), Inches(12.35), Inches(1.6),
                [
                    "原子化 + 数据源解耦 + 降级可观测 + 基线优先",
                    "P0 25 个能力闭环钓鱼 / 勒索 / 异常登录三类告警",
                    "P1 +27 / P2 +31 持续扩展,数据源按依赖序接入",
                ],
                size=14, line_spacing=1.4)

    # link
    add_rect(s, Inches(0.5), Inches(6.05), Inches(12.35), Inches(0.55), NAVY)
    add_text(s, Inches(0.5), Inches(6.05), Inches(12.35), Inches(0.55),
             "配套实施计划:docs/plans/2026-05-24-investigation-atomic-skills-implementation.md",
             size=12, bold=True, color=WHITE, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.5), Inches(6.75), Inches(12.35), Inches(0.4),
             "Q&A · 讨论入口:tech@ouraca.ai",
             size=14, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)


# -------------------------------------------------------------------------
# Driver
# -------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 24
    s01_cover(prs, total)
    s02_agenda(prs, 2, total)
    s03_philosophy(prs, 3, total)
    s04_namespace(prs, 4, total)
    s05_capability_card(prs, 5, total)
    s06_ds_p0(prs, 6, total)
    s07_ds_p1_p2(prs, 7, total)
    s08_overview(prs, 8, total)
    s09_a1_card(prs, 9, total)
    s10_a_others(prs, 10, total)
    s11_b_category(prs, 11, total)
    s12_c_category(prs, 12, total)
    s13_d_category(prs, 13, total)
    s14_e_category(prs, 14, total)
    s15_f_category(prs, 15, total)
    s16_g_category(prs, 16, total)
    s17_matrix(prs, 17, total)
    s18_phishing(prs, 18, total)
    s19_ransomware(prs, 19, total)
    s20_exfil_mining(prs, 20, total)
    s21_bruteforce_lateral(prs, 21, total)
    s22_p0(prs, 22, total)
    s23_p1_p2(prs, 23, total)
    s24_summary(prs, 24, total)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH, len(prs.slides)


if __name__ == "__main__":
    path, n = build()
    size_kb = path.stat().st_size / 1024
    print(f"OK: wrote {path} · slides={n} · size={size_kb:.1f} KB")
